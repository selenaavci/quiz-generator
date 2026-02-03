from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from text_cleaner import clean_extracted_text

SKIP_PLACEHOLDERS = {
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
}

def load_ppt(path: str) -> str:
    try:
        prs = Presentation(path)
        out = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ptype = shape.placeholder_format.type
                        if ptype in SKIP_PLACEHOLDERS:
                            continue
                    except Exception:
                        pass

                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)

        text = "\n".join(out).strip()
        return clean_extracted_text(text) if text else ""

    except Exception as e:
        raise RuntimeError(f"PPTX okunamadı: {e}")

